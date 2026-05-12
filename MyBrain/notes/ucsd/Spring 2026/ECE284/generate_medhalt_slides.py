"""Generate Med-HALT presentation slides 9-15 as a .pptx file.

Follows the layout spec in Pal_2023_MedHALT_Slides_9_to_15_Layout.md:
- 16:9 ratio
- Font sizes >= 16pt (no smaller)
- Color palette: navy / red / green / amber / grey
- ASCII mockup positions

Output: MedHALT_Slides_9_to_15.pptx (same dir).
Then Javen imports into Google Slides via "File → Import slides".
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ─── Palette (hex → RGB) ──────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x3A, 0x6B)
RED = RGBColor(0xCC, 0x4C, 0x54)
GREEN = RGBColor(0x2D, 0x8A, 0x4E)
AMBER = RGBColor(0xD9, 0xA7, 0x3C)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_RED = RGBColor(0xFC, 0xE8, 0xE9)  # very faint red bg
LIGHT_GREEN = RGBColor(0xE5, 0xF3, 0xE9)  # very faint green bg


# ─── Helper functions ─────────────────────────────────────────────
def add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    bold: bool = False,
    color: RGBColor = BLACK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    italic: bool = False,
    font_name: str = "Helvetica",
):
    """Add a text box at (left, top) Inches with width, height Inches.
    font_size in pt. Returns the textbox shape.
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_multi_line(
    slide,
    lines: list,  # list of (text, font_size, bold, color, italic) tuples
    left: float,
    top: float,
    width: float,
    height: float,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    """Add a text box with multiple paragraphs, each with its own formatting."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    for i, line_spec in enumerate(lines):
        text, font_size, bold, color, italic = line_spec
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Helvetica"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=1.0):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_title(slide, title: str, color=NAVY, font_size=36):
    """Standard title at top center, 36 pt bold navy."""
    add_text_box(
        slide,
        title,
        left=0.5,
        top=0.3,
        width=12.3,
        height=0.7,
        font_size=font_size,
        bold=True,
        color=color,
        align=PP_ALIGN.LEFT,
    )
    # Add a thin horizontal divider line under title
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.05), Inches(12.8), Inches(1.05))
    line.line.color.rgb = NAVY
    line.line.width = Pt(1.5)


def new_slide(prs):
    """Create a blank slide and return it."""
    blank_layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(blank_layout)


# ─── Generate presentation ────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════════════
    # Slide 9 — Three Surprises from the Leaderboard
    # ═══════════════════════════════════════════════════════════════
    slide = new_slide(prs)
    add_title(slide, "Three Surprises from the Leaderboard")

    # Left: simplified leaderboard table
    # Build it as a real PowerPoint table
    table_left, table_top = 0.5, 1.5
    table_w, table_h = 7.5, 5.5
    rows, cols = 6, 3
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(table_left), Inches(table_top), Inches(table_w), Inches(table_h)
    )
    tbl = tbl_shape.table
    # column widths
    tbl.columns[0].width = Inches(4.0)
    tbl.columns[1].width = Inches(1.75)
    tbl.columns[2].width = Inches(1.75)

    # Headers
    headers = ["Model", "RHT Avg", "MHT Avg"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = h
        run.font.name = "Helvetica"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = NAVY
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFA)

    # Data rows
    data = [
        ("LLaMA-2 70B Base", "72.33%  ⭐", "8.04%"),
        ("LLaMA-2 70B Chat", "11.26%  ❌", "13.05%"),
        ("Text-Davinci-003", "54.46%", "19.75%"),
        ("Falcon 40B", "59.09%", "30.36%  ⭐"),
        ("GPT-3.5 Turbo", "44.48%", "19.96%"),
    ]
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = val
            run.font.name = "Helvetica"
            run.font.size = Pt(22)
            # Highlight ⭐ and ❌ rows
            if "⭐" in val or "❌" in val:
                run.font.bold = True
                if "⭐" in val:
                    run.font.color.rgb = GREEN
                else:
                    run.font.color.rgb = RED
            else:
                run.font.color.rgb = BLACK

    # Right: RHT / MHT recap
    add_multi_line(
        slide,
        [
            ("RHT = Reasoning Hallucination Test", 22, True, NAVY, False),
            ("", 14, False, BLACK, False),
            ("  • 3 tasks: FCT / NOTA / FQT", 20, False, BLACK, False),
            ("  • Tests robustness vs misleading prompts", 20, False, BLACK, False),
            ("", 16, False, BLACK, False),
            ("MHT = Memory Hallucination Test", 22, True, NAVY, False),
            ("", 14, False, BLACK, False),
            ("  • 4 tasks: PubMed retrieval", 20, False, BLACK, False),
            ("  • Tests honest \"Unknown\" on fake PMIDs", 20, False, BLACK, False),
        ],
        left=8.4,
        top=1.5,
        width=4.6,
        height=5.5,
    )

    # ═══════════════════════════════════════════════════════════════
    # Slide 10 — Surprise 1: Instruction-Tuning Paradox
    # ═══════════════════════════════════════════════════════════════
    slide = new_slide(prs)
    add_title(slide, "Surprise 1 — Instruction-Tuning Paradox 🥇")

    # Left: main comparison
    add_multi_line(
        slide,
        [
            ("LLaMA-2 70B Base", 26, True, BLACK, False),
            ("", 12, False, BLACK, False),
            ("RHT  72.33%  ⭐", 54, True, GREEN, False),
            ("", 10, False, BLACK, False),
            ("↓  + Instruction Tuning + RLHF", 22, True, RED, False),
            ("", 10, False, BLACK, False),
            ("LLaMA-2 70B Chat", 26, True, BLACK, False),
            ("", 12, False, BLACK, False),
            ("RHT  11.26%  ❌", 54, True, RED, False),
            ("", 14, False, BLACK, False),
            ("── Same architecture, 61-point drop ──", 22, True, NAVY, False),
        ],
        left=0.6,
        top=1.3,
        width=6.8,
        height=5.5,
        align=PP_ALIGN.CENTER,
    )

    # Right: Base vs Chat detail
    add_multi_line(
        slide,
        [
            ("Base vs Chat — what's the difference?", 24, True, NAVY, False),
            ("", 12, False, BLACK, False),
            ("Base model", 22, True, BLACK, False),
            ("  = Stops after pre-training", 20, False, BLACK, False),
            ("  → Knowledgeable but can't chat / follow instructions", 20, False, GREY, True),
            ("", 12, False, BLACK, False),
            ("Chat model", 22, True, BLACK, False),
            ("  = Base + SFT + RLHF", 20, False, BLACK, False),
            ("    • SFT = Supervised Fine-Tuning", 17, False, GREY, True),
            ("       (teaches conversation format)", 16, False, GREY, True),
            ("    • RLHF = Reinforcement Learning", 17, False, GREY, True),
            ("       from Human Feedback", 17, False, GREY, True),
            ("       (tunes for human preference)", 16, False, GREY, True),
            ("  → ChatGPT, Claude, Gemini are Chat models", 20, False, GREY, True),
            ("", 14, False, BLACK, False),
            ("⚡ Same architecture + same pre-training", 20, True, AMBER, False),
            ("   — differs only in the final 2 training steps", 20, True, AMBER, False),
        ],
        left=7.8,
        top=1.3,
        width=5.2,
        height=5.5,
    )

    # Bottom: paper quote
    add_text_box(
        slide,
        "📄 Paper §6.1: \"detrimental effect on model's ability to control hallucination after instruction tuning and RLHF.\"",
        left=0.6,
        top=6.8,
        width=12.2,
        height=0.5,
        font_size=16,
        italic=True,
        color=GREY,
    )

    # ═══════════════════════════════════════════════════════════════
    # Slide 11 — Surprise 2: Reasoning ≠ Memory
    # ═══════════════════════════════════════════════════════════════
    slide = new_slide(prs)
    add_title(slide, "Surprise 2 — Reasoning ≠ Memory 🥈")

    # Subtitle
    add_text_box(
        slide,
        "Reasoning and Memory are independent skills.",
        left=0.5,
        top=1.3,
        width=12.3,
        height=0.6,
        font_size=24,
        bold=True,
        color=BLACK,
        align=PP_ALIGN.LEFT,
    )

    # Main table
    rows2, cols2 = 3, 3
    tbl2 = slide.shapes.add_table(
        rows2, cols2, Inches(1.5), Inches(2.4), Inches(10.3), Inches(3.0)
    ).table
    tbl2.columns[0].width = Inches(3.5)
    tbl2.columns[1].width = Inches(3.4)
    tbl2.columns[2].width = Inches(3.4)

    headers2 = ["", "RHT (Reasoning)", "MHT (Memory)"]
    for j, h in enumerate(headers2):
        cell = tbl2.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = "Helvetica"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = NAVY
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFA)

    table_data = [
        ("LLaMA-2 70B", "72.33%  ⭐\n(strong reasoning)", "8.04%\n(weak memory)"),
        ("Falcon 40B", "59.09%\n(mid reasoning)", "30.36%  ⭐\n(strong memory)"),
    ]
    for i, row in enumerate(table_data, start=1):
        for j, val in enumerate(row):
            cell = tbl2.cell(i, j)
            cell.text = ""
            # Split on newline
            for ln_idx, line in enumerate(val.split("\n")):
                if ln_idx == 0:
                    p = cell.text_frame.paragraphs[0]
                else:
                    p = cell.text_frame.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = line
                run.font.name = "Helvetica"
                if j == 0:
                    run.font.size = Pt(22)
                    run.font.bold = True
                    run.font.color.rgb = BLACK
                elif "⭐" in line:
                    run.font.size = Pt(32) if ln_idx == 0 else Pt(18)
                    run.font.bold = True
                    run.font.color.rgb = GREEN if ln_idx == 0 else GREY
                else:
                    run.font.size = Pt(28) if ln_idx == 0 else Pt(18)
                    run.font.bold = ln_idx == 0
                    run.font.color.rgb = BLACK if ln_idx == 0 else GREY

    # Take-away
    add_text_box(
        slide,
        "💡 Modular design: one model for reasoning, another for facts.",
        left=0.6,
        top=5.7,
        width=12.2,
        height=0.5,
        font_size=22,
        italic=True,
        color=NAVY,
    )

    # Side note caveat
    add_text_box(
        slide,
        "— Side note: open-source beats closed-source on RHT, but mostly because closed models are RLHF'd (Surprise 1 again).",
        left=0.6,
        top=6.5,
        width=12.2,
        height=0.7,
        font_size=16,
        italic=True,
        color=GREY,
    )

    # ═══════════════════════════════════════════════════════════════
    # Slide 12 — Surprise 3: FCT 全军覆没
    # ═══════════════════════════════════════════════════════════════
    slide = new_slide(prs)
    add_title(slide, "Surprise 3 — FCT: Where Every Model Fails 🥉")

    # Red banner: No model passes 50%
    add_rect(slide, 0.5, 1.3, 12.3, 0.8, LIGHT_RED, line_color=RED, line_width=1.5)
    add_text_box(
        slide,
        "🚨  No model passes 50%",
        left=0.5,
        top=1.4,
        width=12.3,
        height=0.6,
        font_size=32,
        bold=True,
        color=RED,
        align=PP_ALIGN.CENTER,
    )

    # FCT ranking table
    add_multi_line(
        slide,
        [
            ("False Confidence Test (FCT) accuracy:", 22, True, NAVY, False),
            ("", 12, False, BLACK, False),
            ("LLaMA-2 70B Base       42.21%   (best)", 24, True, GREEN, False),
            ("GPT-3.5 Turbo          34.15%", 22, False, BLACK, False),
            ("Falcon 40B Base        18.66%", 22, False, BLACK, False),
            ("Text-Davinci-003       16.76%", 22, False, BLACK, False),
            ("LLaMA-2 70B Chat       13.34%", 22, False, BLACK, False),
            ("Falcon 40B Instruct     1.11%   (worst)", 24, True, RED, False),
        ],
        left=0.6,
        top=2.4,
        width=7.5,
        height=4.0,
    )

    # Punch line
    add_multi_line(
        slide,
        [
            ("💡 LLM amplifies the user's", 24, True, NAVY, False),
            ("    wrong assumption", 24, True, NAVY, False),
            ("", 10, False, BLACK, False),
            ("(confirmation bias amplifier)", 18, True, AMBER, True),
        ],
        left=8.3,
        top=2.8,
        width=4.8,
        height=2.5,
    )

    # Lyme recall footnote
    add_text_box(
        slide,
        "— Recall Yixian's Lyme example earlier — that was FCT in action.",
        left=0.6,
        top=6.7,
        width=12.2,
        height=0.5,
        font_size=16,
        italic=True,
        color=GREY,
    )

    # ═══════════════════════════════════════════════════════════════
    # Slide 13 — Clinical Implications
    # ═══════════════════════════════════════════════════════════════
    slide = new_slide(prs)
    add_title(slide, "Can we deploy this in hospitals?")

    # 4 row scorecard
    scorecard = [
        ("❌", RED, "Autonomous diagnosis / treatment", "Best model only 72.33% RHT, FCT under 50%", LIGHT_RED),
        ("❌", RED, "RLHF chat models for medical Q&A", "Instruction-Tuning Paradox (Chat: 11.26%)", LIGHT_RED),
        ("✅", GREEN, "Literature search assist (with verification)", "Falcon 40B 30.36% MHT — limited but useful", LIGHT_GREEN),
        ("✅", GREEN, "Physician second-opinion", "Doctor must be in the loop", LIGHT_GREEN),
    ]

    top_start = 1.4
    row_h = 1.0
    for i, (icon, icon_color, claim, why, bg) in enumerate(scorecard):
        y = top_start + i * (row_h + 0.05)
        # Background rectangle
        add_rect(slide, 0.6, y, 12.1, row_h, bg)
        # Icon (large)
        add_text_box(
            slide,
            icon,
            left=0.8,
            top=y + 0.1,
            width=0.9,
            height=row_h - 0.2,
            font_size=44,
            bold=True,
            color=icon_color,
            align=PP_ALIGN.CENTER,
        )
        # Main claim
        add_text_box(
            slide,
            claim,
            left=1.9,
            top=y + 0.1,
            width=10.5,
            height=0.55,
            font_size=24,
            bold=True,
            color=BLACK,
        )
        # Why
        add_text_box(
            slide,
            why,
            left=1.9,
            top=y + 0.55,
            width=10.5,
            height=0.4,
            font_size=18,
            italic=True,
            color=GREY,
        )

    # Punch line at bottom
    add_text_box(
        slide,
        "🩺  Augmentation, not Automation",
        left=0.5,
        top=6.8,
        width=12.3,
        height=0.6,
        font_size=32,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    # ═══════════════════════════════════════════════════════════════
    # Slide 14 — Limitations & Open Questions
    # ═══════════════════════════════════════════════════════════════
    slide = new_slide(prs)
    add_title(slide, "Limitations & Open Questions")

    # Left column: paper acknowledged
    add_multi_line(
        slide,
        [
            ("📄 Paper acknowledges", 24, True, NAVY, False),
            ("", 14, False, BLACK, False),
            ("1. MCQ ≠ free-form clinical Q&A", 20, False, BLACK, False),
            ("", 8, False, BLACK, False),
            ("2. Only 7 tasks (no planning / multimodal)", 20, False, BLACK, False),
            ("", 8, False, BLACK, False),
            ("3. Prompt brittleness unsolved", 20, False, BLACK, False),
            ("", 8, False, BLACK, False),
            ("4. GPT-4 / Claude not tested", 20, False, BLACK, False),
            ("    (publication timing)", 18, False, GREY, True),
        ],
        left=0.6,
        top=1.4,
        width=6.0,
        height=5.5,
    )

    # Vertical separator line
    sep = slide.shapes.add_connector(
        1, Inches(6.8), Inches(1.5), Inches(6.8), Inches(6.8)
    )
    sep.line.color.rgb = LIGHT_GREY
    sep.line.width = Pt(1.5)

    # Right column: my open questions
    add_multi_line(
        slide,
        [
            ("💡 My open questions", 24, True, NAVY, False),
            ("", 14, False, BLACK, False),
            ("5. IT-Paradox: paper shows", 20, False, BLACK, False),
            ('   "what" not "why" — needs ablation', 20, False, GREY, True),
            ("", 12, False, BLACK, False),
            ("6. Med-PaLM 2 NOT Med-HALT-ed", 20, False, BLACK, False),
            ("    — clear research gap", 20, False, GREY, True),
            ("", 12, False, BLACK, False),
            ("7. Pointwise Score → RLHF reward", 20, False, BLACK, False),
            ("    function?  (train honest LLMs)", 20, False, GREY, True),
        ],
        left=7.1,
        top=1.4,
        width=6.0,
        height=5.5,
    )

    # ═══════════════════════════════════════════════════════════════
    # Slide 15 — Discussion
    # ═══════════════════════════════════════════════════════════════
    slide = new_slide(prs)
    add_title(slide, "Discussion", font_size=44)

    # Q1 (Clinical)
    add_multi_line(
        slide,
        [
            ("Q1   (Clinical)", 26, True, NAVY, False),
            ("", 10, False, BLACK, False),
            (
                "If you were the CMO of a hospital deciding whether to deploy",
                22,
                False,
                BLACK,
                False,
            ),
            (
                "an LLM for physician second-opinion in radiology, which finding",
                22,
                False,
                BLACK,
                False,
            ),
            (
                "from this paper would worry you the MOST — the 16.76% FCT score,",
                22,
                False,
                BLACK,
                False,
            ),
            (
                "the LLaMA-Chat collapse, or prompt brittleness?",
                22,
                False,
                BLACK,
                False,
            ),
            ("Why that one specifically?", 22, True, RED, False),
        ],
        left=0.6,
        top=1.4,
        width=12.2,
        height=2.6,
    )

    # Horizontal separator
    sep2 = slide.shapes.add_connector(
        1, Inches(0.6), Inches(4.2), Inches(12.8), Inches(4.2)
    )
    sep2.line.color.rgb = LIGHT_GREY
    sep2.line.width = Pt(1.5)

    # Q2 (Technical)
    add_multi_line(
        slide,
        [
            ("Q2   (Technical)", 26, True, NAVY, False),
            ("", 10, False, BLACK, False),
            (
                "RLHF makes the chat model WORSE at hallucination control.",
                22,
                False,
                BLACK,
                False,
            ),
            (
                'If you had a research budget — would you (a) fix RLHF with a',
                22,
                False,
                BLACK,
                False,
            ),
            (
                '"honesty reward," (b) skip RLHF entirely for medical models,',
                22,
                False,
                BLACK,
                False,
            ),
            (
                "or (c) do 2-stage training?",
                22,
                False,
                BLACK,
                False,
            ),
            ("Defend your choice.", 22, True, RED, False),
        ],
        left=0.6,
        top=4.4,
        width=12.2,
        height=2.6,
    )

    # ─── Save ─────────────────────────────────────────────────────
    out_path = "MedHALT_Slides_9_to_15.pptx"
    prs.save(out_path)
    print(f"✓ Generated: {out_path}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Size: 16:9 ({prs.slide_width / 914400:.2f}\" × {prs.slide_height / 914400:.2f}\")")


if __name__ == "__main__":
    main()
